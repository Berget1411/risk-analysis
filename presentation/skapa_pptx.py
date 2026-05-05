"""Genererar presentation.pptx för ME1316-projektet. ~10 min presentation.

Slides: Titel + 12 innehållsslides med inline-tabeller och grafer.
Kör: uv run python presentation/skapa_pptx.py
"""
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = Path(__file__).parent
FIGURE_DIR = OUTPUT_DIR / "figures"
FIGURE_DIR.mkdir(exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT = RGBColor(0x00, 0x6B, 0x8F)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)


# === Helpers ===

def add_title_slide(title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = title_text
    for para in title.text_frame.paragraphs:
        para.font.size = Pt(40)
        para.font.bold = True
        para.font.color.rgb = DARK

    subtitle = slide.placeholders[1]
    subtitle.text = subtitle_text
    for para in subtitle.text_frame.paragraphs:
        para.font.size = Pt(22)
        para.font.color.rgb = ACCENT

    return slide


def add_slide(title_text, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = title_text
    for para in title.text_frame.paragraphs:
        para.font.size = Pt(32)
        para.font.bold = True
        para.font.color.rgb = DARK

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    if subtitle:
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = ACCENT
        p.space_after = Pt(12)
        start_new = True
    else:
        start_new = False

    for i, bullet in enumerate(bullets):
        if i == 0 and not start_new:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet.startswith("  - "):
            p.text = bullet.strip("- ").strip()
            p.level = 1
            p.font.size = Pt(18)
        else:
            p.text = bullet.lstrip("- ").strip()
            p.level = 0
            p.font.size = Pt(20)

        p.font.color.rgb = DARK
        p.space_after = Pt(6)

    return slide


def add_table_slide(title_text, headers, rows, col_widths=None, subtitle=None):
    """Add a slide with a simple table."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.italic = True
        p2.font.color.rgb = ACCENT

    # Table
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    left = Inches(0.8)
    top = Inches(1.6) if not subtitle else Inches(1.9)
    width = Inches(11.5)
    height = Inches(0.4) * n_rows

    if col_widths is None:
        col_widths = [width / n_cols] * n_cols

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = int(w)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(15)
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT

    return slide


def add_figure_slide(title_text, fig_path, subtitle=None):
    """Add a slide with a figure image."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.italic = True
        p2.font.color.rgb = ACCENT

    # Image centered
    slide.shapes.add_picture(
        str(fig_path),
        Inches(1.5), Inches(1.3),
        width=Inches(10), height=Inches(5.8),
    )
    return slide


# === Generate figures ===

def make_claim_frequency_bar():
    """Horizontal bar chart: skadefrekvens per verksamhet + geografi (side by side)."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))

    # Verksamhet
    labels_v = ["VVS", "Takarbeten", "Byggföretag", "Övr. specialist", "Grävn. & Schakt.", "Elektriker", "Målare"]
    values_v = [0.0316, 0.0250, 0.0221, 0.0214, 0.0189, 0.0154, 0.0143]
    colors_v = ["#d32f2f" if v > 0.025 else "#1976d2" if v < 0.016 else "#546e7a" for v in values_v]

    ax1.barh(labels_v, values_v, color=colors_v)
    ax1.set_xlabel("Skador per exponerat år")
    ax1.set_title("Per verksamhet")
    ax1.axvline(0.0214, ls="--", color="gray", alpha=0.5, label="Portföljsnitt")
    ax1.legend(fontsize=9)
    ax1.invert_yaxis()

    # Geografi
    labels_g = ["Storstad", "Mellanstorstad", "Landsbyggd", "Småstad"]
    values_g = [0.0261, 0.0215, 0.0179, 0.0135]
    colors_g = ["#d32f2f", "#546e7a", "#546e7a", "#1976d2"]

    ax2.barh(labels_g, values_g, color=colors_g)
    ax2.set_xlabel("Skador per exponerat år")
    ax2.set_title("Per geografi")
    ax2.axvline(0.0214, ls="--", color="gray", alpha=0.5, label="Portföljsnitt")
    ax2.legend(fontsize=9)
    ax2.invert_yaxis()

    plt.tight_layout()
    path = FIGURE_DIR / "skadefrekvens_segment.png"
    plt.savefig(path, dpi=150, bbox_inches="tight")
    plt.close()
    return path


def make_decile_bar():
    """Vertical bar chart: skadefrekvens per omsättningsdecil (1-10)."""
    fig, ax = plt.subplots(figsize=(10, 5))

    deciles = list(range(1, 11))
    values = [0.0072, 0.0100, 0.0130, 0.0150, 0.0180, 0.0200, 0.0230, 0.0270, 0.0320, 0.0400]

    # Color gradient blue → red
    cmap = plt.cm.RdYlBu_r
    norm = plt.Normalize(vmin=min(values), vmax=max(values))
    colors = [cmap(norm(v)) for v in values]

    ax.bar(deciles, values, color=colors, edgecolor="white", linewidth=0.5)
    ax.axhline(0.0214, ls="--", color="gray", alpha=0.7, label="Portföljsnitt (0,0214)")
    ax.set_xlabel("Omsättningsdecil (1 = lägst, 10 = högst)")
    ax.set_ylabel("Skadefrekvens")
    ax.set_title("Skadefrekvens per omsättningsdecil — 5,6× spridning")
    ax.set_xticks(deciles)
    ax.legend(fontsize=10)
    ax.set_ylim(0, 0.045)

    # Annotate spread
    ax.annotate("5,6×", xy=(9.5, 0.041), fontsize=14, fontweight="bold", color="#d32f2f",
                ha="center")

    plt.tight_layout()
    path = FIGURE_DIR / "omsattning_decil.png"
    plt.savefig(path, dpi=150, bbox_inches="tight")
    plt.close()
    return path


def make_model_comparison_bar():
    """Simple grouped bar showing GLM vs XGBoost key metrics on 2025 test."""
    fig, ax = plt.subplots(figsize=(10, 5))

    metrics = ["Poisson\ndeviance", "Portföljfel\n(%)", "Predikterade\nskador"]
    glm_vals = [41889.2, 1.10, 5581]
    xgb_vals = [41855.7, 1.22, 5587]

    x = np.arange(len(metrics))
    w = 0.3

    bars1 = ax.bar(x - w/2, glm_vals, w, label="GLM M2", color="#1976d2")
    bars2 = ax.bar(x + w/2, xgb_vals, w, label="XGBoost", color="#ff8f00")

    # Add observed line for predicted claims
    ax.axhline(5520, xmin=0.65, xmax=0.98, ls="--", color="green", alpha=0.7, label="Observerat (5 520)")

    ax.set_xticks(x)
    ax.set_xticklabels(metrics)
    ax.legend()
    ax.set_title("Testportfölj 2025: GLM M2 vs XGBoost")

    # Annotate bars
    for bar in bars1:
        h = bar.get_height()
        ax.annotate(f'{h:,.1f}' if h > 100 else f'{h:.2f}%',
                    xy=(bar.get_x() + bar.get_width()/2, h),
                    xytext=(0, 3), textcoords="offset points",
                    ha='center', fontsize=9)
    for bar in bars2:
        h = bar.get_height()
        ax.annotate(f'{h:,.1f}' if h > 100 else f'{h:.2f}%',
                    xy=(bar.get_x() + bar.get_width()/2, h),
                    xytext=(0, 3), textcoords="offset points",
                    ha='center', fontsize=9)

    plt.tight_layout()
    path = FIGURE_DIR / "modelljaemfoerelse.png"
    plt.savefig(path, dpi=150, bbox_inches="tight")
    plt.close()
    return path


def make_rate_ratio_forest():
    """Forest plot of rate ratios with confidence intervals."""
    fig, ax = plt.subplots(figsize=(10, 6))

    labels = [
        "VVS", "Takarbeten", "Övr. specialist",
        "Grävn. & Schakt.", "Elektriker", "Målare",
        "",  # separator
        "Storstad", "Mellanstorstad", "Småstad",
        "",  # separator
        "Omsättning (×2)",
    ]
    rr =     [1.432, 1.135, 0.970, 0.855, 0.698, 0.637, None, 1.461, 1.203, 0.757, None, 1.358]
    ci_low = [1.372, 1.067, 0.932, 0.808, 0.659, 0.601, None, 1.387, 1.140, 0.711, None, 1.345]
    ci_hi =  [1.494, 1.207, 1.010, 0.905, 0.738, 0.675, None, 1.540, 1.270, 0.806, None, 1.371]

    y_pos = []
    plot_labels = []
    plot_rr = []
    plot_ci_low = []
    plot_ci_hi = []
    y = 0
    for i, label in enumerate(labels):
        if label == "":
            y += 0.5
            continue
        y_pos.append(y)
        plot_labels.append(label)
        plot_rr.append(rr[i])
        plot_ci_low.append(ci_low[i])
        plot_ci_hi.append(ci_hi[i])
        y += 1

    y_pos = np.array(y_pos)
    plot_rr = np.array(plot_rr)
    plot_ci_low = np.array(plot_ci_low)
    plot_ci_hi = np.array(plot_ci_hi)

    # Colors based on above/below 1
    colors = ["#d32f2f" if r > 1 else "#1976d2" for r in plot_rr]

    ax.axvline(1.0, color="gray", ls="--", alpha=0.5)
    ax.errorbar(plot_rr, y_pos, xerr=[plot_rr - plot_ci_low, plot_ci_hi - plot_rr],
                fmt="none", ecolor="gray", capsize=3, alpha=0.7)
    ax.scatter(plot_rr, y_pos, c=colors, s=80, zorder=5)

    ax.set_yticks(y_pos)
    ax.set_yticklabels(plot_labels)
    ax.set_xlabel("Rate Ratio (ref: Byggföretag / Landsbyggd)")
    ax.set_title("Rate Ratios med 95% KI — GLM M2")
    ax.invert_yaxis()
    ax.grid(axis="x", alpha=0.3)

    # Section labels
    ax.text(0.55, 0, "Verksamhet", fontsize=10, fontstyle="italic", color="gray",
            transform=ax.get_yaxis_transform())
    ax.text(0.55, 6.5, "Geografi", fontsize=10, fontstyle="italic", color="gray",
            transform=ax.get_yaxis_transform())
    ax.text(0.55, 9.5, "Storlek", fontsize=10, fontstyle="italic", color="gray",
            transform=ax.get_yaxis_transform())

    plt.tight_layout()
    path = FIGURE_DIR / "rate_ratios_forest.png"
    plt.savefig(path, dpi=150, bbox_inches="tight")
    plt.close()
    return path


# Generate all figures
fig_segments = make_claim_frequency_bar()
fig_decile = make_decile_bar()
fig_forest = make_rate_ratio_forest()
fig_comparison = make_model_comparison_bar()

print(f"Figurer sparade i {FIGURE_DIR}/")


# === SLIDES ===

# 1. Titel
add_title_slide(
    "Skadefrekvensanalys\nEntreprenadförsäkring",
    "ME1316 · Poisson-GLM vs XGBoost · 10 min"
)

# 2. Bakgrund och Syfte
add_slide("Bakgrund och Syfte", [
    "Entreprenadförsäkring: försäkring för byggprojekt och anläggninsarbeten",
    "Premie = frekvens × kostnad — vi analyserar frekvenskomponenten",
    "",
    "Syfte:",
    "  - Identifiera vilka faktorer som driver skadefrekvensen",
    "  - Prediktera skadeantal för en ny portfölj (2025)",
    "  - Jämföra tolkbar modell (GLM) mot maskininlärning (XGBoost)",
])

# 3. Frågeställning och Metod
add_slide("Frågeställning och Metod", [
    "Tre forskningsfrågor:",
    "  - Vilka faktorer driver skadefrekvensen?",
    "  - Kan vi prediktera 2025 med god träffsäkerhet?",
    "  - Ger XGBoost mervärde jämfört med GLM?",
    "",
    "Tidsbaserad validering (undviker look-ahead bias):",
    "  - Träna: 2021–2023",
    "  - Validera: 2024 (modellval och hyperparametrar)",
    "  - Testa: 2025 (slutlig utvärdering, aldrig sedd under utveckling)",
])

# 4. Dataöverblick (tabell)
add_table_slide(
    "Dataöverblick",
    ["", "Träning", "Test"],
    [
        ["Period", "2021–2024", "2025"],
        ["Antal rader", "1 033 386", "291 649"],
        ["Antal skador", "19 730", "5 520"],
        ["Noll-skador", "98,1%", "98,1%"],
        ["Skadefrekvens", "0,0214", "0,0212"],
    ],
    subtitle="7 verksamhetstyper · 4 geografiska områden · Omsättning (kontinuerlig)",
)

# 5. Deskriptiva fynd — segmentfigur
add_figure_slide(
    "Tre ratingfaktorer sticker ut",
    fig_segments,
    subtitle="Skadefrekvens varierar 2–3× mellan segment. VVS + Storstad = högst risk.",
)

# 6. Deskriptiva fynd — Omsättning decil
add_figure_slide(
    "Omsättning: starkaste gradienten",
    fig_decile,
    subtitle="5,6× spridning mellan lägsta och högsta decil. Motiverar log(Omsättning) i modellen.",
)

# 7. Poisson-GLM: Modellval (AIC-tabell)
add_table_slide(
    "Poisson-GLM: Modellval",
    ["Modell", "Beskrivning", "AIC", "ΔAIC vs M0"],
    [
        ["M0", "Enbart intercept", "141 086", "—"],
        ["M1", "Verksamhet + Geografi", "139 723", "−1 363"],
        ["M2", "M1 + log(Omsättning)", "136 971", "−4 115"],
        ["M3", "M2 + År", "136 969", "−4 117"],
    ],
    subtitle="M2 vald som slutmodell. Dispersion 0,986 → Poisson-antagandet håller. Årseffekt negligerbar.",
)

# 8. Rate Ratios (forest plot)
add_figure_slide(
    "Rate Ratios — multiplikatorer på skadefrekvens",
    fig_forest,
    subtitle="Allt annat lika. Röd = högre risk, blå = lägre risk vs referens.",
)

# 9. Rate Ratios tabell
add_table_slide(
    "Nyckeltal: Rate Ratios",
    ["Faktor", "Rate Ratio", "95% KI", "Tolkning"],
    [
        ["VVS (vs Byggföretag)", "1,43", "[1,37 – 1,49]", "+43% skadefrekvens"],
        ["Storstad (vs Landsbyggd)", "1,46", "[1,39 – 1,54]", "+46% skadefrekvens"],
        ["Omsättning ×2", "1,36", "[1,35 – 1,37]", "+36% per fördubbling"],
        ["Målare (vs Byggföretag)", "0,64", "[0,60 – 0,68]", "−36% skadefrekvens"],
        ["Småstad (vs Landsbyggd)", "0,76", "[0,71 – 0,81]", "−24% skadefrekvens"],
    ],
    subtitle="Snäva KI → skillnaderna är statistiskt säkra",
)

# 10. XGBoost: Challenger
add_slide("XGBoost: Challenger-modell", [
    "Gradient boosting med Poisson-objektiv, samma 3 variabler",
    "Bästa konfiguration: max_depth = 3, learning_rate = 0,10, 232 träd",
    "  - Grunda träd → signalen är additiv, precis som GLM antar",
    "  - Djupare träd (djup 5, 8) ger sämre resultat",
    "",
    "Feature importance (gain):",
    "  - log(Omsättning): 20,04",
    "  - GeografisktOmråde: 8,26",
    "  - Verksamhet: 6,46",
    "",
    "Samma tre drivare, samma rangordning → GLM missar inget",
])

# 11. Modelljämförelse (tabell)
add_table_slide(
    "GLM vs XGBoost — Testportfölj 2025",
    ["Mått", "GLM M2", "XGBoost", "Skillnad"],
    [
        ["Poisson deviance", "41 889", "41 856", "−0,08%"],
        ["RMSE", "0,1374", "0,1374", "≈ 0"],
        ["Portföljfel", "+1,10%", "+1,22%", "GLM bättre"],
        ["Pred. skador", "5 581", "5 587", "Obs: 5 520"],
    ],
    subtitle="Validering 2024 visade samma mönster (−0,21% deviance). Ingen meningsfull skillnad.",
)

# 12. Osäkerhet
add_slide("Osäkerhet i prediktioner", [
    "Portföljnivå:",
    "  - Totalprognos: 5 581 skador, 95% KI [5 503 – 5 659]",
    "  - Relativ osäkerhet: 2,8% → portföljprognosen är tillförlitlig",
    "  - Observerat (5 520) ligger inom KI ✓",
    "",
    "Radnivå: relativ osäkerhet varierar 0,05 – 0,19",
    "",
    "Drivare av hög osäkerhet:",
    "  - Liten träningscell (ovanliga verksamhet × geografi-kombinationer)",
    "  - Extrema omsättningsnivåer (långt från centrum)",
    "  - Kort exponeringstid (lågt förväntat antal)",
])

# 13. Slutsats och Rekommendation
add_slide("Slutsats och Rekommendation", [
    "GLM predikterar 5 581 skador vs 5 520 observerade (+1,1% fel)",
    "  - 3× spridning från lägsta till högsta riskprofil",
    "",
    "Rekommendation:",
    "  - Differentierad premie: Verksamhet × Geografi × Omsättning",
    "  - GLM som huvudmodell (tolkbar, transparent, stabil)",
    "  - XGBoost som parallell challenger (larmar vid modelldrift)",
    "  - Manuell granskning för ovanliga kundprofiler med hög osäkerhet",
    "",
    "Begränsningar:",
    "  - Samband, ej orsak. Frekvens, ej kostnad.",
    "  - Vidare: kostnadsmodell, interaktioner, finare variabler, årlig uppföljning",
])

# === Save ===
out_path = OUTPUT_DIR / "risk-presentation.pptx"
prs.save(str(out_path))
print(f"\nSparad: {out_path}")
